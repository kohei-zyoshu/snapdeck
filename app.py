"""
パシャッと — 手書きメモ・ホワイトボード → スライド自動変換
"""

import os, io, json, base64, re, copy
from datetime import datetime, timezone, timedelta

JST = timezone(timedelta(hours=9))

import streamlit as st
from PIL import Image, ImageChops, ImageEnhance, ImageOps

# ─── ページ設定 ────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="パシャッと",
    page_icon="📸",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ─── CSS（20代女性向け・グラデーションバイオレット） ────────────────────────
st.markdown("""
<style>
/* ── Streamlit UI を非表示 ── */
header[data-testid="stHeader"] { display: none !important; }
#MainMenu { display: none !important; }
footer { display: none !important; }
[data-testid="stToolbar"] { display: none !important; }
[data-testid="stDecoration"] { display: none !important; }
a[href*="github.com"] { display: none !important; }
a[href*="streamlit.io"] { display: none !important; }
.viewerBadge_container__r5tak { display: none !important; }
.stActionButton { display: none !important; }

/* ── モバイルタップ最適化（効果音・ハイライト抑制） ── */
* { -webkit-tap-highlight-color: transparent; }
button, a { touch-action: manipulation; -webkit-appearance: none; }

/* ── 全体背景（クリーンホワイト） ── */
.stApp { background: #FFFFFF; }
.main .block-container { padding: 0 1rem 5rem; max-width: 560px; }

/* ── 固定ヘッダーバー（グラデーション） ── */
.top-bar {
    position: sticky;
    top: 0;
    z-index: 999;
    background: linear-gradient(135deg, #6D28D9 0%, #A855F7 55%, #EC4899 100%);
    padding: 0.6rem 1.2rem;
    margin: 0 -1rem 1.6rem;
    display: flex;
    align-items: center;
    gap: 0.6rem;
    box-shadow: 0 3px 16px rgba(109,40,217,0.30);
}
.top-bar-logo {
    font-size: 1.55rem;
    font-weight: 900;
    color: #ffffff;
    letter-spacing: -0.01em;
    line-height: 1;
}
.top-bar-logo em { color: #FDE8FF; font-style: normal; }
.top-bar-tagline {
    font-size: 0.88rem;
    color: rgba(255,255,255,0.80);
    font-weight: 500;
    line-height: 1.35;
}

/* ── ステップラベル ── */
.step {
    display: flex;
    align-items: center;
    gap: 0.75rem;
    font-size: 1.35rem;
    font-weight: 800;
    color: #1E1B4B;
    margin: 2rem 0 0.6rem;
}
.snum {
    background: linear-gradient(135deg, #7C3AED, #A855F7);
    color: #fff;
    border-radius: 50%;
    width: 40px; height: 40px;
    line-height: 40px;
    text-align: center;
    font-size: 1.1rem;
    font-weight: 900;
    flex-shrink: 0;
    box-shadow: 0 3px 10px rgba(124,58,237,0.35);
}

/* ── ヒント・説明文 ── */
.hint {
    color: #3B0764;
    font-size: 1.05rem;
    margin: 0 0 1rem;
    line-height: 1.8;
    background: #F5F3FF;
    border-left: 4px solid #7C3AED;
    border-radius: 0 8px 8px 0;
    padding: 0.75rem 1rem;
}

/* ── ボタン全般 ── */
.stButton > button {
    font-size: 1.2rem !important;
    min-height: 66px !important;
    border-radius: 18px !important;
    font-weight: 800 !important;
    width: 100% !important;
    letter-spacing: 0.02em !important;
    transition: opacity 0.12s, transform 0.1s !important;
    cursor: pointer !important;
}
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #7C3AED 0%, #A855F7 50%, #EC4899 100%) !important;
    color: white !important;
    border: none !important;
    box-shadow: 0 5px 20px rgba(124,58,237,0.40) !important;
}
.stButton > button[kind="primary"]:hover  { opacity: 0.88 !important; }
.stButton > button[kind="primary"]:active { transform: scale(0.97) !important; }
.stButton > button[kind="secondary"] {
    background: #F5F3FF !important;
    color: #4C1D95 !important;
    border: 2px solid #C4B5FD !important;
}
.stButton > button:disabled { opacity: 0.35 !important; cursor: not-allowed !important; }

/* ── ダウンロードボタン（メイン：グラデーション） ── */
.stDownloadButton > button {
    font-size: 1.2rem !important;
    min-height: 66px !important;
    border-radius: 18px !important;
    font-weight: 800 !important;
    width: 100% !important;
    border: none !important;
    letter-spacing: 0.02em !important;
    transition: opacity 0.12s, transform 0.1s !important;
    cursor: pointer !important;
}
.stDownloadButton > button:active { transform: scale(0.97) !important; }

/* ── 完了ボックス ── */
.done-box {
    background: #FAF5FF;
    border: 2px solid #A855F7;
    border-radius: 14px;
    padding: 1.2rem 1.4rem;
    font-size: 1.1rem;
    color: #4C1D95;
    font-weight: 700;
    margin: 1.2rem 0;
    line-height: 1.8;
}

/* ── ファイルアップローダー ── */
[data-testid="stFileUploaderDropzone"] {
    min-height: 150px;
    border-radius: 14px !important;
    border: 2.5px dashed #C4B5FD !important;
    background: #FAFAFF !important;
}
[data-testid="stFileUploaderDropzoneInstructions"] span,
[data-testid="stFileUploaderDropzoneInstructions"] small { display: none !important; }
[data-testid="stFileUploaderDropzoneInstructions"]::before {
    content: "ドラッグまたはボタンで選択";
    display: block;
    font-size: 1.0rem;
    color: #5B21B6;
    margin-bottom: 0.5rem;
    text-align: center;
}
[data-testid="stFileUploaderDropzoneInstructions"]::after {
    content: "JPG・PNG・WebP・PDF（最大20MB）";
    display: block;
    font-size: 0.85rem;
    color: #8B5CF6;
    text-align: center;
    margin-top: 0.3rem;
}
[data-testid="stFileUploaderDropzone"] button {
    font-size: 0 !important;
    min-height: 50px !important;
    border-radius: 12px !important;
    padding: 0 1.5rem !important;
}
[data-testid="stFileUploaderDropzone"] button::after {
    content: "ファイルを選ぶ";
    font-size: 1.0rem;
    font-weight: 700;
}

/* ── その他 ── */
.stToggle label { font-size: 1.05rem !important; font-weight: 700 !important; }
.stExpander summary p { font-size: 1.05rem !important; font-weight: 600 !important; }
hr { margin: 1.6rem 0 !important; border-color: #EDE9FE !important; border-width: 1.5px !important; }
.stAlert p { font-size: 1.0rem !important; line-height: 1.7 !important; }
.stSpinner p { font-size: 1.05rem !important; }
</style>
""", unsafe_allow_html=True)


# ─── バックエンド関数 ──────────────────────────────────────────────────────────

def get_api_key() -> str:
    """APIキーを Secrets / 環境変数 / 入力欄から取得"""
    try:
        key = st.secrets["ANTHROPIC_API_KEY"]
        if key:
            return str(key).strip()
    except Exception:
        pass
    env_key = os.environ.get("ANTHROPIC_API_KEY", "").strip()
    if env_key:
        return env_key
    return st.session_state.get("api_key_input", "").strip()


def get_pdf_info(file_bytes: bytes) -> int:
    """PDFのページ数を返す"""
    import fitz  # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    count = doc.page_count
    doc.close()
    return count


def pdf_page_to_pil(file_bytes: bytes, page_num: int = 0) -> Image.Image:
    """PDFの指定ページを PIL Image（RGB）に変換（2倍解像度）"""
    import fitz
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    page = doc[page_num]
    mat = fitz.Matrix(2.0, 2.0)        # 2×ズームで解像度を確保
    pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
    doc.close()
    return img


def preprocess_image(img: Image.Image,
                     do_trim: bool = True) -> Image.Image:
    """EXIF回転 → OCR最適化補正（自動）→ 余白カット

    補正方針:
    1. 照明むらが大きい場合のみ Retinex 正規化（影・逆光を自動検出）
    2. autocontrast でヒストグラムを常に最大伸張
    3. シャープネス強化でテキスト輪郭を際立たせる
    """
    import numpy as np
    from PIL import ImageFilter

    try:
        img = ImageOps.exif_transpose(img)
    except Exception:
        pass
    if img.mode != "RGB":
        img = img.convert("RGB")

    # ── 照明むら自動検出 ──
    # ブラー画像の輝度標準偏差が大きいほど照明が不均一
    radius   = max(40, min(img.width, img.height) // 6)
    blur_pil = img.filter(ImageFilter.GaussianBlur(radius=radius))
    blur_l   = np.array(blur_pil.convert("L"), dtype=np.float32)
    if float(blur_l.std()) > 22:      # 照明むらが目立つ → Retinex で均一化
        arr  = np.array(img, dtype=np.float32)
        blur = np.array(blur_pil, dtype=np.float32)
        norm = arr / (blur / 160.0 + 0.5)
        img  = Image.fromarray(np.clip(norm, 0, 255).astype(np.uint8))

    # ── ヒストグラム自動伸張（常に適用） ──
    img = ImageOps.autocontrast(img, cutoff=1)

    # ── シャープネス強化 ──
    img = ImageEnhance.Sharpness(img).enhance(1.5)

    # ── 余白カット ──
    if do_trim:
        bg   = Image.new("RGB", img.size, (255, 255, 255))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()
        if bbox:
            margin = 30
            img = img.crop((
                max(0,          bbox[0] - margin),
                max(0,          bbox[1] - margin),
                min(img.width,  bbox[2] + margin),
                min(img.height, bbox[3] + margin),
            ))
    return img


def pil_to_base64(img: Image.Image, max_px: int = 1280) -> tuple[str, str]:
    """PIL Image → OCR用 base64（preprocess 済み画像をそのままエンコード）
    max_px=1280: 手書きOCRに必要十分な解像度。2048と比べ画像トークン数が
    約半分になりAPI応答が3〜5秒速くなる。"""
    if img.width > max_px or img.height > max_px:
        img.thumbnail((max_px, max_px), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=95)
    return base64.standard_b64encode(buf.getvalue()).decode("utf-8"), "image/jpeg"


@st.cache_data(show_spinner=False)
def cached_process_image(file_bytes: bytes,
                         do_trim: bool) -> tuple[str, str, bool, bytes]:
    """
    画像のバイト列を受け取り、前処理 → base64変換 → プレビュー用バイト列を返す。
    同じ入力なら Streamlit がキャッシュするため、リランのたびに再処理されない。
    戻り値: (img_data_b64, media_type, is_portrait, preview_jpeg_bytes)
    """
    pil = Image.open(io.BytesIO(file_bytes))
    if pil.mode != "RGB":
        pil = pil.convert("RGB")
    disp = preprocess_image(pil, do_trim=do_trim)
    is_portrait = disp.height > disp.width
    img_data, media_type = pil_to_base64(disp)
    # プレビュー用にJPEGバイト列として保持（PIL Imageはキャッシュ不可）
    preview_buf = io.BytesIO()
    disp.save(preview_buf, format="JPEG", quality=88)
    return img_data, media_type, is_portrait, preview_buf.getvalue()


# 色帯パレット（アイテムごとに割り当てる 12 色）
_BAND_COLORS = [
    (231, 76,  60),   # red
    (52,  152, 219),  # blue
    (39,  174, 96),   # green
    (155, 89,  182),  # purple
    (230, 126, 34),   # orange
    (32,  178, 170),  # teal
    (236, 72,  153),  # pink
    (99,  110, 250),  # indigo
    (243, 156, 18),   # yellow
    (76,  201, 240),  # cyan
    (67,  97,  238),  # royal blue
    (247, 37,  133),  # hot pink
]

def item_color_hex(global_idx: int) -> str:
    r, g, b = _BAND_COLORS[global_idx % len(_BAND_COLORS)]
    return f"#{r:02x}{g:02x}{b:02x}"




EXTRACTION_PROMPT = """あなたは優秀なOCRエンジンです。
この画像（ホワイトボード・手書きメモ・付箋・資料など）の内容を正確に読み取り、プレゼンテーション用に整理してください。

【解析手順】
STEP1: 画像全体を俯瞰して「1列構成か2列構成か」「表があるか」を確認する。
STEP2: 上から下・左から右の順にすべての文字と表を正確に読み取る。
STEP3: 出現順に blocks 配列へ記録する。
  ・文章・箇条書き・見出しのかたまり → type:"section"
  ・罫線・格子状の行列構造 → type:"table"
  ・2列構成なら左側を column:1、右側を column:2 とする

【返却JSON形式】
"title":"タイトル","blocks":[{"type":"section","heading":"見出し（なければ空文字）","column":1,"items":[{"text":"大見出し","type":"heading","shape":"none","color":"black","bold":true,"x_pct":25,"y_pct":10},{"text":"箇条書き","type":"bullet","shape":"none","color":"black","bold":false,"x_pct":25,"y_pct":25}]},{"type":"table","headers":["列1","列2"],"rows":[["値1","値2"]]}],"lines":[{"x1_pct":10,"y1_pct":50,"x2_pct":30,"y2_pct":50,"color":"black","style":"solid","type":"line"}]}

【必須ルール】
- title はスライドヘッダー専用（blocks の items に含めない）
- blocks は画像内の出現順（上から下）に並べる
- section の column: 左側・1列レイアウトは 1、右側コンテンツは 2
- type: heading / bullet / text / arrow
- shape: rect（四角囲み）/ ellipse（丸囲み）/ none（囲みなし）
- color: ペン・インクの実際の色を正確に識別する（black / red / blue / green / orange / purple / pink / gray / brown）
  ★重要: 黒いペンは "black"、青いペン・マーカーは "blue"、赤いペンは "red" — 実際の色を忠実に設定すること
  ★例: 青いマーカーで書かれた日付 → color:"blue" / 黒ボールペンで書かれたメモ → color:"black"
- x_pct: そのテキスト・付箋の中心が画像左端から何%の位置か（0〜100の整数）
- y_pct: そのテキスト・付箋の中心が画像上端から何%の位置か（0〜100の整数）
- bg_color: 付箋・カードの背景色を実際の色で正確に識別する（yellow/pink/blue/green/orange/purple/white/none）
  ★黄色と黄緑・緑の区別: 純粋な黄色（レモン色）→ "yellow" / 黄緑・ライムグリーン・緑系 → "green"
  ★重要: ホワイトボード・紙に直接書かれた文字（付箋の外のテキスト）は必ず "none"
  ★重要: 付箋（ポストイット）の上に書かれた文字のみ実際の付箋色（yellow/pink/blue等）を設定
  ★例: ホワイトボードに書いた列ヘッダー・日付・タイトルは bg_color:"none"
- table の headers: 見出し行がない場合は []
- lines: ホワイトボード上に描かれた線・矢印（テキストや付箋の外にある線のみ）を配列で記録
  type: line（直線）/ arrow（片方向矢印、x2/y2 側が矢じり）/ double_arrow（両方向矢印）
  style: solid（実線）/ dashed（破線）
  color: 線のペン色（black / blue / red 等）
  x1_pct/y1_pct: 始点の位置（0〜100）、x2_pct/y2_pct: 終点の位置（0〜100）
  線がない場合は lines:[] とする
- 画像内のすべての文字を漏れなく・正確に読む（推測・省略・合体禁止）
- 読み取れなかった・自信のない文字は [?] で示す（例: "田中[?]"、"合計[?]円"）
- 必ず { で始まる JSONのみ返す（説明文・コードブロック・前置き不要）"""


def _repair_truncated_json(text: str) -> str:
    """末尾が途切れた JSON を修復する。
    スタックで開いているブラケット/ブレースを追跡し、末尾に閉じ記号を補完する。
    途中で切断された文字列リテラルも閉じる。
    """
    stack: list[str] = []
    in_string    = False
    escape_next  = False

    for ch in text:
        if escape_next:
            escape_next = False
            continue
        if ch == "\\" and in_string:
            escape_next = True
            continue
        if ch == '"':
            in_string = not in_string
            continue
        if in_string:
            continue
        if ch in "{[":
            stack.append("}" if ch == "{" else "]")
        elif ch in "}]":
            if stack and stack[-1] == ch:
                stack.pop()

    suffix = ""
    if in_string:
        suffix += '"'                        # 開いたままの文字列を閉じる
    suffix += "".join(reversed(stack))       # 開いたままの {} [] を閉じる
    return text + suffix


def _parse_json_response(raw: str) -> dict:
    """AI応答テキストを堅牢にJSONパースして返す。
    失敗時は末尾切断とみなし修復を試みる。
    """
    text = re.sub(r"```(?:json)?\s*", "", raw)
    text = re.sub(r"```", "", text).strip()
    start = text.find("{")
    end   = text.rfind("}")
    if start == -1:
        raise ValueError(f"応答にJSONが見つかりません（先頭: {raw[:200]}）")

    # 末尾 } がない場合は全体を修復対象にする
    body = text[start:(end + 1) if end > start else len(text)]
    body = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", body)

    # ① そのままパース
    try:
        return json.loads(body)
    except json.JSONDecodeError:
        pass

    # ② 末尾切断の修復を試みる
    repaired = _repair_truncated_json(body)
    try:
        return json.loads(repaired)
    except json.JSONDecodeError as e:
        raise ValueError(f"JSON解析エラー: {e}\n先頭: {body[:200]}")


def _normalize_blocks(data: dict) -> dict:
    """旧フォーマット（sections + tables）を blocks[] に統一し、
    各ブロックの null 値を修復する。"""
    # ── 新フォーマット（blocks あり）の正規化 ──
    if data.get("blocks"):
        for b in data["blocks"]:
            if b.get("type") == "section":
                if b.get("items") is None:
                    b["items"] = []
                for item in b["items"]:
                    for key in ("x_pct", "y_pct"):
                        v = item.get(key)
                        if v is not None:
                            item[key] = max(0, min(100, int(v)))
        # lines の正規化
        if data.get("lines") is None:
            data["lines"] = []
        for ln in data["lines"]:
            for key in ("x1_pct", "y1_pct", "x2_pct", "y2_pct"):
                v = ln.get(key)
                if v is not None:
                    ln[key] = max(0, min(100, int(v)))
        return data

    # ── 旧フォーマット（sections + tables）→ blocks に変換 ──
    blocks: list[dict] = []
    for sec in (data.get("sections") or []):
        if sec.get("items") is None:
            sec["items"] = []
        blocks.append({"type": "section", **sec})
    # 旧 items フラット形式の互換
    if not blocks and data.get("items"):
        blocks.append({
            "type": "section", "heading": "", "column": 1,
            "items": [
                {"text": it.get("text", ""), "type": it.get("type", "text"),
                 "shape": it.get("shape", "none"), "color": it.get("color", "black"),
                 "bold": it.get("bold", False)}
                for it in (data["items"] or [])
            ],
        })
    for tbl in (data.get("tables") or []):
        blocks.append({"type": "table", **tbl})
    data["blocks"] = blocks
    return data


def analyze_with_claude(img_data: str, media_type: str, api_key: str,
                        model: str = "claude-sonnet-4-6") -> dict:
    """AIで画像を解析し、blocks[] 形式の構造化データを返す。

    精度向上のための施策:
    - temperature=0: JSON出力の再現性を高める
    - assistant prefill: 出力を即座に { で始めさせ、前置きを排除
    - _parse_json_response: 堅牢な3ステップJSONクリーニング
    """
    import anthropic
    client = anthropic.Anthropic(api_key=api_key)
    message = client.messages.create(
        model=model,
        max_tokens=8192,        # 大型ホワイトボード対応（旧3000は切断頻発）
        temperature=0,          # 決定論的出力でJSON崩れを防ぐ
        messages=[
            {"role": "user", "content": [
                {"type": "image", "source": {
                    "type": "base64", "media_type": media_type, "data": img_data
                }},
                {"type": "text", "text": EXTRACTION_PROMPT},
            ]},
        ]
    )
    raw = message.content[0].text.strip()

    data = _parse_json_response(raw)
    data = _normalize_blocks(data)

    # ── UI表示用 elements リストをセクションブロックから構築 ──
    elements = []
    for block in data.get("blocks", []):
        if block.get("type") != "section":
            continue
        for it in (block.get("items") or []):
            idx = len(elements)
            elements.append({
                "id":      f"el_{idx:03d}",
                "type":    it.get("type", "text"),
                "content": it.get("text", ""),
                "shape":   it.get("shape", "none"),
                "style":   {"color": it.get("color", "black"),
                             "bold":  it.get("bold", False)},
                "confidence": 1.0,
            })
    data["elements"] = elements
    data.setdefault("structure", {"type": "list", "groups": [
        {"label": "内容", "items": [e["id"] for e in elements]}
    ]})
    return data


VERIFY_PROMPT = """この画像から以下のテキストが自動抽出されました。
元の画像と1件ずつ照合し、識字精度を評価してください。

【抽出テキスト一覧】
{items_json}

【返却形式】
[{{"idx":0,"confidence":"high","correction":null}},{{"idx":1,"confidence":"medium","correction":"修正後テキスト"}},...]

confidence:
- "high"  … 正確に読み取れている
- "medium" … やや不確か、別の読み方の可能性あり
- "low"   … 誤読の可能性が高い

correction: 修正後テキスト（high の場合は null、[?] を含む場合は補完案を記載）
JSONのみ返す（説明文不要）。"""


def verify_extraction(img_data: str, media_type: str, api_key: str,
                      data: dict) -> dict:
    """抽出結果を元画像と照合して識字精度を評価する（常に Haiku を使用）。

    戻り値: {(bi, ii): {"confidence": "high"|"medium"|"low", "correction": str|None}}
    """
    import anthropic

    # section ブロックのアイテムを番号付きリストに展開
    items: list[dict] = []
    for bi, block in enumerate(data.get("blocks") or []):
        if block.get("type") != "section":
            continue
        for ii, item in enumerate(block.get("items") or []):
            items.append({"idx": len(items), "bi": bi, "ii": ii,
                          "text": item.get("text", "")})

    if not items:
        return {}

    items_for_prompt = [{"idx": it["idx"], "text": it["text"]} for it in items]
    prompt = VERIFY_PROMPT.format(
        items_json=json.dumps(items_for_prompt, ensure_ascii=False))

    client  = anthropic.Anthropic(api_key=api_key)
    message = client.messages.create(
        model="claude-haiku-4-5-20251001",   # 速さ重視・コスト最小
        max_tokens=2000,
        temperature=0,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {
                "type": "base64", "media_type": media_type, "data": img_data
            }},
            {"type": "text", "text": prompt},
        ]}],
    )
    raw = message.content[0].text.strip()

    # JSON 配列を抽出
    text = re.sub(r"```(?:json)?\s*", "", raw)
    text = re.sub(r"```", "", text).strip()
    ls   = text.find("[")
    le   = text.rfind("]")
    if ls == -1 or le == -1:
        return {}
    try:
        results: list[dict] = json.loads(text[ls:le + 1])
    except json.JSONDecodeError:
        return {}

    # idx → (bi, ii) のマッピング
    idx_map = {it["idx"]: (it["bi"], it["ii"]) for it in items}
    out: dict = {}
    for r in results:
        key = idx_map.get(r.get("idx"))
        if key is not None:
            out[key] = {
                "confidence": r.get("confidence", "medium"),
                "correction": r.get("correction"),
            }
    return out


REREAD_PROMPT = """この画像から以下のテキストが読み取られましたが、一部の文字が不確かです。
画像を注意深く見て、各テキストを正確に読み直してください。

【再読み取りリスト】
{items_json}

【返却形式】
[{{"idx":0,"text":"正確なテキスト"}},{{"idx":1,"text":"正確なテキスト"}},...]

- 元の読み取り結果より画像内の実際の文字を優先する
- [?] の部分は読み取れた文字で補完する（読めなければそのまま残す）
- JSONのみ返す（説明文不要）"""


def reread_uncertain(img_data: str, media_type: str, api_key: str,
                     data: dict, verification: dict,
                     model: str = "claude-sonnet-4-6") -> dict:
    """🟡/🔴 の項目のみ高精度モデルで再読み取りする。

    コスト最小化のポイント:
    - 不確かな箇所のテキストリストのみ送信（出力トークン数が極小）
    - 全体の再解析ではなく差分のみ修正
    戻り値: {(bi, ii): "修正後テキスト"}
    """
    import anthropic

    # 🟡/🔴 の項目だけを対象にする
    targets: list[dict] = []
    for (bi, ii), v in verification.items():
        if v.get("confidence") not in ("medium", "low"):
            continue
        blks = data.get("blocks") or []
        if bi < len(blks):
            items = blks[bi].get("items") or []
            if ii < len(items):
                targets.append({
                    "idx": len(targets), "bi": bi, "ii": ii,
                    "text": items[ii].get("text", ""),
                    "confidence": v["confidence"],
                })

    if not targets:
        return {}

    items_for_prompt = [
        {"idx": t["idx"], "text": t["text"],
         "note": "low" if t["confidence"] == "low" else "medium"}
        for t in targets
    ]
    prompt = REREAD_PROMPT.format(
        items_json=json.dumps(items_for_prompt, ensure_ascii=False))

    client  = anthropic.Anthropic(api_key=api_key)
    message = client.messages.create(
        model=model,
        max_tokens=1000,    # 出力は短いテキストリストのみ
        temperature=0,
        messages=[{"role": "user", "content": [
            {"type": "image", "source": {
                "type": "base64", "media_type": media_type, "data": img_data
            }},
            {"type": "text", "text": prompt},
        ]}],
    )
    raw  = message.content[0].text.strip()
    text = re.sub(r"```(?:json)?\s*", "", raw)
    text = re.sub(r"```", "", text).strip()
    ls   = text.find("[")
    le   = text.rfind("]")
    if ls == -1 or le == -1:
        return {}
    try:
        results: list[dict] = json.loads(text[ls:le + 1])
    except json.JSONDecodeError:
        return {}

    idx_map = {t["idx"]: (t["bi"], t["ii"]) for t in targets}
    return {idx_map[r["idx"]]: r["text"]
            for r in results if r.get("idx") in idx_map and r.get("text")}


def generate_pptx(data: dict, is_portrait: bool = False) -> bytes:
    """解析データからパワーポイントを生成（縦横自動対応）"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ─ カラーパレット ─
    C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)
    C_NAVY   = RGBColor(0x1B, 0x28, 0x38)
    C_ACCENT = RGBColor(0x00, 0xB4, 0xD8)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_MUTED  = RGBColor(0x94, 0xA3, 0xB8)
    C_LIGHT  = RGBColor(0xCB, 0xD5, 0xE1)

    # ─ 手書き色 → RGBColor ─
    TEXT_COLORS = {
        "black":  RGBColor(0xF0, 0xF4, 0xF8),
        "white":  RGBColor(0xFF, 0xFF, 0xFF),
        "red":    RGBColor(0xFF, 0x6B, 0x6B),
        "blue":   RGBColor(0x64, 0xB5, 0xF6),
        "green":  RGBColor(0x81, 0xC7, 0x84),
        "yellow": RGBColor(0xFF, 0xF1, 0x76),
        "orange": RGBColor(0xFF, 0xB7, 0x4D),
        "purple": RGBColor(0xCE, 0x93, 0xD8),
        "pink":   RGBColor(0xF4, 0x8F, 0xB1),
        "gray":   RGBColor(0xB0, 0xBE, 0xC5),
        "brown":  RGBColor(0xBC, 0xAA, 0xA4),
    }

    def resolve_color(name: str) -> RGBColor:
        return TEXT_COLORS.get((name or "black").lower(), C_LIGHT)

    # ─ 縦横でスライドサイズを切り替え ─
    if is_portrait:
        SLIDE_W, SLIDE_H = 7.5, 10.83   # A4縦
    else:
        SLIDE_W, SLIDE_H = 13.33, 7.5   # 16:9横

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    def bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def rect(slide, x, y, w, h, fc, lc=None):
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fc
        if lc:
            sh.line.color.rgb = lc
        else:
            sh.line.fill.background()

    def txt(slide, text, x, y, w, h, size=13, bold=False, color=None,
            align=PP_ALIGN.LEFT, wrap=True, underline=False, italic=False):
        color = color or C_LIGHT
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.underline = underline
        r.font.italic    = italic
        r.font.color.rgb = color
        r.font.name      = "Noto Sans JP"

    def shape_txt(slide, shape_type, text, x, y, w, h,
                  size=13, bold=False, color=None, underline=False):
        """テキスト付きオートシェイプ（rect / ellipse）を追加"""
        SHAPE_ID = {"rect": 1, "ellipse": 9, "rounded_rect": 5}
        sid = SHAPE_ID.get(shape_type, 1)
        color = color or C_LIGHT
        sh = slide.shapes.add_shape(sid,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
        # 塗りつぶし：ネイビー + テキスト色のボーダー
        sh.fill.solid()
        sh.fill.fore_color.rgb = C_NAVY
        sh.line.color.rgb = color
        sh.line.width = Pt(1.5)
        tf = sh.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text           = text
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.underline = underline
        r.font.color.rgb = color
        r.font.name      = "Noto Sans JP"

    def add_table_shape(slide, tbl_data, x, y, max_w, max_h):
        headers = tbl_data.get("headers", [])
        rows    = tbl_data.get("rows", [])
        if not headers and not rows:
            return
        n_rows = len(rows) + (1 if headers else 0)
        n_cols = max(len(headers), max((len(r) for r in rows), default=1))
        if n_rows == 0 or n_cols == 0:
            return
        col_w = min(max_w / n_cols, 2.5)
        row_h = min(max_h / n_rows, 0.45)
        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(x), Inches(y),
            Inches(col_w * n_cols), Inches(row_h * n_rows)
        ).table
        if headers:
            for ci, h in enumerate(headers[:n_cols]):
                cell = tbl.cell(0, ci)
                cell.text = str(h)
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_ACCENT
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold      = True
                        run.font.size      = Pt(11)
                        run.font.color.rgb = C_DARK
        for ri, row in enumerate(rows):
            tr = ri + (1 if headers else 0)
            for ci, val in enumerate(row[:n_cols]):
                cell = tbl.cell(tr, ci)
                cell.text = str(val)
                bg_col = RGBColor(0x1E, 0x30, 0x40) if ri % 2 == 0 else C_NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_col
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size      = Pt(10)
                        run.font.color.rgb = C_LIGHT

    # ── スライド作成 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C_DARK)
    rect(s, 0, 0, SLIDE_W, 0.07, C_ACCENT)

    title    = data.get("title", "ホワイトボードメモ")
    ts_label = datetime.now(JST).strftime("%Y年%m月%d日 %H:%M")
    CONTENT_X = 0.4
    TITLE_W   = SLIDE_W - 0.8

    txt(s, title, CONTENT_X, 0.10, TITLE_W, 0.65,
        size=20 if is_portrait else 24, bold=True, color=C_WHITE)
    txt(s, f"変換日時：{ts_label}", CONTENT_X, 0.73, TITLE_W, 0.30,
        size=9, color=C_MUTED)
    rect(s, 0.3, 1.08, SLIDE_W - 0.6, SLIDE_H - 1.55, C_NAVY)

    # ── コンテンツエリア定義 ──
    # タイトルバー下からフッター上までを「コンテンツエリア」とする
    CX = 0.35              # 左マージン（インチ）
    CY = 1.10              # コンテンツ開始Y（インチ）
    CW = SLIDE_W - 0.70    # コンテンツ幅
    CH = SLIDE_H - 1.55    # コンテンツ高さ（フッター分を引く）

    blocks = data.get("blocks") or []

    # ── ブロック分類 ──
    sec_blocks = [b for b in blocks if b.get("type") == "section"]
    col1_secs  = [b for b in sec_blocks if b.get("column", 1) != 2]
    col2_secs  = [b for b in sec_blocks if b.get("column", 1) == 2]
    has_two_cols = bool(col2_secs) and not is_portrait

    def count_rows(secs: list) -> int:
        """セクションリストの総行数（見出し行 + アイテム行）"""
        n = 0
        for sec in secs:
            if sec.get("heading"):
                n += 1
            n += len(sec.get("items") or [])
        return max(n, 1)

    # ── フォントサイズをセクション行数から逆算 ──
    n_sec_rows = (max(count_rows(col1_secs), count_rows(col2_secs))
                  if has_two_cols else count_rows(sec_blocks))
    EL_H    = max(0.24, min(0.44, CH / max(n_sec_rows, 1)))
    base_pt = max(9, min(16, int(EL_H * 72 * 0.50)))
    head_pt = min(base_pt + 3, 20)
    sec_pt  = min(base_pt + 5, 22)

    def render_item_at(item: dict, x_start: float, width: float, y: float) -> None:
        """1アイテムを指定座標に描画する（y は呼び出し元が管理）"""
        etype      = item.get("type", "text")
        content    = item.get("text", "")
        shape_type = item.get("shape", "none")
        color      = resolve_color(item.get("color", "black"))
        bold       = item.get("bold", False) or etype == "heading"
        if shape_type in ("rect", "ellipse"):
            shape_txt(s, shape_type, content,
                      x_start + 0.05, y, width - 0.1, EL_H,
                      size=base_pt, bold=bold, color=color)
        elif etype == "heading":
            txt(s, content, x_start, y, width, EL_H,
                size=head_pt, bold=True, color=color)
        elif etype == "arrow":
            txt(s, f"→ {content}", x_start + 0.1, y, width - 0.1, EL_H,
                size=base_pt, color=RGBColor(0xFF, 0xB7, 0x4D), italic=True)
        elif etype == "bullet":
            txt(s, f"・{content}", x_start + 0.1, y, width - 0.1, EL_H,
                size=base_pt, bold=bold, color=color)
        else:
            txt(s, content, x_start + 0.1, y, width - 0.1, EL_H,
                size=base_pt, bold=bold, color=color)

    def render_section_list(secs: list, x_start: float, width: float,
                            start_y: float = CY) -> float:
        """セクションリストを描画し、終端 Y を返す"""
        y = start_y
        for sec in secs:
            heading = sec.get("heading", "").strip()
            if heading:
                txt(s, heading, x_start, y, width, EL_H,
                    size=sec_pt, bold=True, color=C_WHITE, underline=True)
                y += EL_H
            for item in (sec.get("items") or []):
                if y >= CY + CH:
                    break
                render_item_at(item, x_start, width, y)
                y += EL_H
        return y

    if has_two_cols:
        # ── 2列モード：左右セクションを並列描画、表はその後に順番通り ──
        col_w = (CW - 0.30) / 2
        render_section_list(col1_secs, CX,                col_w)
        render_section_list(col2_secs, CX + col_w + 0.30, col_w)
        # 表ブロックを2列コンテンツ下に続けて配置
        tbl_y = CY + n_sec_rows * EL_H + 0.1
        for b in blocks:
            if b.get("type") != "table":
                continue
            if tbl_y >= CY + CH:
                break
            tbl_h = min(1.8, CY + CH - tbl_y)
            add_table_shape(s, b, CX, tbl_y, max_w=CW, max_h=tbl_h)
            tbl_y += tbl_h + 0.15
    else:
        # ── 1列モード：blocks を出現順に上から描画 ──
        y = CY
        for b in blocks:
            if y >= CY + CH:
                break
            if b.get("type") == "section":
                y = render_section_list([b], CX, CW, start_y=y)
            elif b.get("type") == "table":
                tbl_h = min(1.8, CY + CH - y)
                if tbl_h > 0.3:
                    add_table_shape(s, b, CX, y + 0.05, max_w=CW, max_h=tbl_h)
                    y += tbl_h + 0.15

    # ─ フッター ─
    FOOTER_Y = SLIDE_H - 0.35
    rect(s, 0, FOOTER_Y, SLIDE_W, 0.35, C_NAVY)
    txt(s, "パシャッと  |  自動生成", 0.4, FOOTER_Y + 0.02, SLIDE_W - 0.5, 0.30,
        size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── ポストイット SVG パレット（セクション単位で色を割り当て）──
_STICKY_PALETTE = [
    ("#FFF9C4", "#F9A825"),  # yellow
    ("#FCE4EC", "#E91E63"),  # pink
    ("#E3F2FD", "#1565C0"),  # blue
    ("#E8F5E9", "#2E7D32"),  # green
    ("#FBE9E7", "#BF360C"),  # orange
    ("#EDE7F6", "#4527A0"),  # purple
    ("#E0F7FA", "#006064"),  # teal
    ("#FFF3E0", "#E65100"),  # amber
    ("#FAFAFA", "#616161"),  # gray
    ("#F3E5F5", "#6A1B9A"),  # deep purple
    ("#E8EAF6", "#283593"),  # indigo
    ("#E0F2F1", "#004D40"),  # deep teal
]

# ── 付箋背景色マップ: bg_color → (fill, stroke) ──
BG_COLOR_MAP: dict[str, tuple[str, str]] = {
    "yellow": ("#FFF9C4", "#F9A825"),
    "pink":   ("#FCE4EC", "#E91E63"),
    "red":    ("#FFEBEE", "#C62828"),
    "blue":   ("#E3F2FD", "#1565C0"),
    "green":  ("#C8E6C9", "#2E7D32"),
    "orange": ("#FFF3E0", "#E65100"),
    "purple": ("#EDE7F6", "#4527A0"),
    "white":  ("#FFFFFF", "#9E9E9E"),
    "gray":   ("#FAFAFA", "#616161"),
    "brown":  ("#EFEBE9", "#4E342E"),
}

# ── ペン色マップ: color → hex ──
PEN_COLOR_MAP: dict[str, str] = {
    "black":  "#1A1A1A",
    "blue":   "#1565C0",
    "red":    "#C62828",
    "green":  "#2E7D32",
    "purple": "#4527A0",
    "pink":   "#AD1457",
    "orange": "#E65100",
    "brown":  "#4E342E",
    "gray":   "#616161",
    "white":  "#FFFFFF",
}


def pen_color_hex(item: dict) -> str:
    """item の color フィールドからペン色 hex を返す（未指定は黒）"""
    return PEN_COLOR_MAP.get((item.get("color") or "black").lower(), "#1A1A1A")


def item_bg_colors(item: dict, palette_idx: int) -> tuple[str, str]:
    """item の bg_color を (fill, stroke) へ変換。未登録はパレットから割り当て"""
    bg = (item.get("bg_color") or "none").lower()
    return BG_COLOR_MAP.get(bg, _STICKY_PALETTE[palette_idx % len(_STICKY_PALETTE)])


def svg_esc(t: str) -> str:
    """SVG テキスト用 XML エスケープ"""
    return (str(t).replace("&", "&amp;").replace("<", "&lt;")
                  .replace(">", "&gt;").replace('"', "&quot;"))


def svg_wrap_text(text: str, max_ch: int = 18) -> list[str]:
    """テキストを max_ch 文字で折り返し、最大4行に収める"""
    t, lines = str(text), []
    while t:
        lines.append(t[:max_ch])
        t = t[max_ch:]
        if len(lines) >= 4:
            if t:
                lines[-1] = lines[-1][:-1] + "…"
            break
    return lines


def generate_svg(data: dict, preview_bytes: bytes | None = None) -> bytes:
    """ホワイトボード・ポストイットの内容をFigmaで編集可能なSVGとして出力する。

    x_pct / y_pct が設定されている場合は元の空間配置を再現（空間モード）。
    ない場合はセクション単位のグリッドレイアウトにフォールバック。
    各付箋は <g id="item-{bi}-{ii}"> で独立したレイヤーになるためFigmaで個別編集可能。
    """
    import io as _io

    # ── キャンバスサイズ（画像のアスペクト比に合わせる）──
    CANVAS_W = 1200
    if preview_bytes:
        try:
            _pil = Image.open(_io.BytesIO(preview_bytes))
            CANVAS_H = int(CANVAS_W * _pil.height / _pil.width)
        except Exception:
            CANVAS_H = 900
    else:
        CANVAS_H = 900

    STICKY_W  = 180
    STICKY_H  = 108
    GAP       = 14
    SEC_PAD   = 20
    SEC_HDR_H = 32
    SEC_GAP   = 28
    PAGE_PAD  = 48
    INNER_W   = CANVAS_W - PAGE_PAD * 2

    # モジュールレベルの関数をローカルエイリアスで参照（クロージャ内で見やすくするため）
    esc        = svg_esc
    text_lines = svg_wrap_text
    pen_color  = pen_color_hex
    item_colors = item_bg_colors

    def svg_sticky(cx: float, cy: float, text: str,
                   fill: str, stroke: str, item_id: str,
                   text_color: str = "#1A1A1A",
                   font_size: int = 12) -> str:
        x       = cx - STICKY_W / 2
        y       = cy - STICKY_H / 2
        max_ch  = max(5, int(STICKY_W / font_size * 1.1))  # 折り返し文字数を絞って大きく見せる
        lines   = text_lines(text, max_ch=max_ch)
        LH      = font_size + 5
        ty0     = cy - (len(lines) - 1) * LH / 2
        rect    = (f'<rect x="{x:.1f}" y="{y:.1f}" width="{STICKY_W}" height="{STICKY_H}"'
                   f' rx="8" fill="{fill}" stroke="{stroke}" stroke-width="1.5"'
                   f' filter="url(#dropshadow)"/>\n')
        txts    = "".join(
            f'<text x="{cx:.1f}" y="{ty0 + i * LH:.1f}" font-size="{font_size}"'
            f' fill="{text_color}" text-anchor="middle" dominant-baseline="central"'
            f' font-family="system-ui,-apple-system,sans-serif">{esc(ln)}</text>\n'
            for i, ln in enumerate(lines)
        )
        return f'<g id="{item_id}">\n{rect}{txts}</g>\n'

    def svg_label(cx: float, cy: float, text: str,
                  item_id: str,
                  text_color: str = "#1E1B4B",
                  font_size: int = 14,
                  bold: bool = True) -> str:
        """付箋なしのテキストラベル（ホワイトボード直書き・列ヘッダー等）"""
        fw     = "700" if bold else "400"
        lines  = text_lines(text, max_ch=20)
        LH     = font_size + 4
        ty0    = cy - (len(lines) - 1) * LH / 2
        txts   = "".join(
            f'<text x="{cx:.1f}" y="{ty0 + i * LH:.1f}" font-size="{font_size}"'
            f' font-weight="{fw}" fill="{text_color}" text-anchor="middle"'
            f' dominant-baseline="central"'
            f' font-family="system-ui,-apple-system,sans-serif">{esc(ln)}</text>\n'
            for i, ln in enumerate(lines)
        )
        return f'<g id="{item_id}">\n{txts}</g>\n'

    def svg_table(tbl: dict, x0: float, y0: float,
                  width: float, tbl_id: str) -> tuple[str, float]:
        headers = tbl.get("headers") or []
        rows    = tbl.get("rows") or []
        n_cols  = max(len(headers), max((len(r) for r in rows), default=0))
        if n_cols == 0:
            return "", 0
        col_w   = width / n_cols
        RH, HH  = 36, 40
        svg     = ""
        cy      = y0
        # 外枠
        total_h = (HH if headers else 0) + len(rows) * RH
        svg += (f'<rect x="{x0:.1f}" y="{y0:.1f}" width="{width:.1f}"'
                f' height="{total_h}" fill="none" stroke="#CBD5E1"'
                f' stroke-width="1" rx="4"/>\n')
        if headers:
            svg += (f'<rect x="{x0:.1f}" y="{cy:.1f}" width="{width:.1f}"'
                    f' height="{HH}" fill="#E3E8EF" rx="4"/>\n')
            for ci, h in enumerate(headers[:n_cols]):
                tx = x0 + ci * col_w + col_w / 2
                svg += (f'<text x="{tx:.1f}" y="{cy + HH/2:.1f}" font-size="13"'
                        f' font-weight="bold" fill="#1E293B" text-anchor="middle"'
                        f' dominant-baseline="central"'
                        f' font-family="system-ui,-apple-system,sans-serif">{esc(h)}</text>\n')
            cy += HH
        for ri, row in enumerate(rows):
            bg = "#F8FAFC" if ri % 2 == 0 else "#FFFFFF"
            svg += (f'<rect x="{x0:.1f}" y="{cy:.1f}" width="{width:.1f}"'
                    f' height="{RH}" fill="{bg}"/>\n')
            for ci, val in enumerate(row[:n_cols]):
                tx = x0 + ci * col_w + col_w / 2
                svg += (f'<text x="{tx:.1f}" y="{cy + RH/2:.1f}" font-size="12"'
                        f' fill="#334155" text-anchor="middle" dominant-baseline="central"'
                        f' font-family="system-ui,-apple-system,sans-serif">{esc(val)}</text>\n')
            cy += RH
        return f'<g id="{tbl_id}">\n{svg}</g>\n', total_h

    # ── アイテムに x_pct / y_pct があるか確認 ──
    blocks    = data.get("blocks") or []
    title     = data.get("title", "")
    has_coord = any(
        item.get("x_pct") is not None and item.get("y_pct") is not None
        for b in blocks if b.get("type") == "section"
        for item in (b.get("items") or [])
    )

    body_svg  = ""
    sec_idx   = 0
    tbl_idx   = 0

    if has_coord:
        # ══ 空間モード：x_pct / y_pct で元の配置を再現 ══
        WB_PAD  = 32
        WB_X    = WB_PAD
        WB_Y    = PAGE_PAD + 56
        WB_W    = CANVAS_W - WB_PAD * 2
        WB_H    = CANVAS_H - WB_Y - WB_PAD
        total_canvas_h = CANVAS_H

        # ── x_pct をクラスタリングして列を均等配置 ──
        # 付箋アイテム（bg_color != "none"）のx_pctだけ使って列数を検出
        sticky_xp = [
            item.get("x_pct", 50)
            for b in blocks if b.get("type") == "section"
            for item in (b.get("items") or [])
            if item.get("x_pct") is not None
            and (item.get("bg_color") or "none").lower() != "none"
        ]
        if not sticky_xp:
            # 付箋がない場合は全アイテムのx_pctを使う
            sticky_xp = [
                item.get("x_pct", 50)
                for b in blocks if b.get("type") == "section"
                for item in (b.get("items") or [])
                if item.get("x_pct") is not None
            ]
        if sticky_xp:
            sorted_xp     = sorted(set(sticky_xp))
            CTHR          = 10   # この%差以内を同じ列とみなす
            col_groups: list[list[float]] = []
            for x in sorted_xp:
                if col_groups and x - col_groups[-1][-1] <= CTHR:
                    col_groups[-1].append(x)
                else:
                    col_groups.append([x])
            n_cols        = max(len(col_groups), 1)
            group_means   = [sum(g) / len(g) for g in col_groups]
            even_x_pcts   = [(i + 0.5) / n_cols * 100 for i in range(n_cols)]

            def snap_x(xp: float) -> float:
                idx = min(range(len(group_means)), key=lambda i: abs(group_means[i] - xp))
                return even_x_pcts[idx]
        else:
            n_cols = 1
            def snap_x(xp: float) -> float:
                return xp

        # ── 列数に応じてポストイットサイズを自動スケール ──
        STICKY_W  = min(200, max(80, int(WB_W / n_cols * 0.78)))
        STICKY_H  = max(80, int(STICKY_W * 0.78))    # 正方形に近い比率（実物のポストイット）
        FONT_SIZE = max(13, int(16 * STICKY_W / 180)) # より大きめのフォント

        # ホワイトボード背景
        body_svg += (f'<rect x="{WB_X}" y="{WB_Y}" width="{WB_W}" height="{WB_H}"'
                     f' rx="12" fill="#FAFAFA" stroke="#E5E7EB" stroke-width="2"/>\n')

        for bi, block in enumerate(blocks):
            if block.get("type") == "section":
                for ii, item in enumerate(block.get("items") or []):
                    xp  = snap_x(item.get("x_pct", 50))
                    yp  = item.get("y_pct", 50)
                    cx  = WB_X + WB_W * xp / 100
                    cy  = WB_Y + WB_H * yp / 100
                    bg  = (item.get("bg_color") or "none").lower()
                    if bg == "none":
                        # ホワイトボード直書き → テキストラベルとして描画
                        body_svg += svg_label(cx, cy, item.get("text", ""),
                                              item_id=f"item-{bi}-{ii}",
                                              text_color=pen_color(item),
                                              font_size=FONT_SIZE + 2,
                                              bold=item.get("type") in ("heading",))
                    else:
                        # 付箋 → 端クランプして付箋として描画
                        cx  = max(WB_X + STICKY_W/2 + 4,
                                  min(WB_X + WB_W - STICKY_W/2 - 4, cx))
                        cy  = max(WB_Y + STICKY_H/2 + 4,
                                  min(WB_Y + WB_H - STICKY_H/2 - 4, cy))
                        fill, stroke = item_colors(item, sec_idx)
                        body_svg += svg_sticky(cx, cy, item.get("text", ""),
                                               fill, stroke, f"item-{bi}-{ii}",
                                               text_color=pen_color(item),
                                               font_size=FONT_SIZE)
                sec_idx += 1
            elif block.get("type") == "table":
                # 表はホワイトボードの下に配置
                s, h = svg_table(block, WB_X, WB_Y + WB_H + SEC_GAP,
                                 WB_W, f"table-{tbl_idx}")
                body_svg += s
                total_canvas_h = WB_Y + WB_H + SEC_GAP + h + WB_PAD
                tbl_idx += 1

        # ── 線・矢印の描画（付箋より下のレイヤーに挿入）──
        lines_svg = ""
        for li, ln in enumerate(data.get("lines") or []):
            lx1 = WB_X + WB_W * snap_x(ln.get("x1_pct", 0)) / 100
            ly1 = WB_Y + WB_H * ln.get("y1_pct", 0) / 100
            lx2 = WB_X + WB_W * snap_x(ln.get("x2_pct", 100)) / 100
            ly2 = WB_Y + WB_H * ln.get("y2_pct", 0) / 100
            lhex   = PEN_COLOR_MAP.get((ln.get("color") or "black").lower(), "#1A1A1A")
            ltype  = (ln.get("type") or "line").lower()
            lstyle = (ln.get("style") or "solid").lower()

            dash_attr  = ' stroke-dasharray="8,4"' if lstyle == "dashed" else ""
            marker_id  = f"ah-{lhex.lstrip('#')}"
            end_attr   = f' marker-end="url(#{marker_id})"'    if ltype in ("arrow", "double_arrow") else ""
            start_attr = f' marker-start="url(#{marker_id})"'  if ltype == "double_arrow" else ""

            lines_svg += (
                f'<line id="line-{li}" x1="{lx1:.1f}" y1="{ly1:.1f}"'
                f' x2="{lx2:.1f}" y2="{ly2:.1f}"'
                f' stroke="{lhex}" stroke-width="2"{dash_attr}'
                f'{start_attr}{end_attr} stroke-linecap="round"/>\n'
            )
        # 線をホワイトボード背景の直後・付箋の直前に挿入
        body_svg = body_svg.replace(
            f'<rect x="{WB_X}" y="{WB_Y}"',
            f'<g id="lines">\n{lines_svg}</g>\n<rect x="{WB_X}" y="{WB_Y}"',
            1,
        )
    else:
        # ══ グリッドモード（フォールバック）：セクション単位で整列 ══
        total_canvas_h = PAGE_PAD + 56

        def grid_section(sec: dict, x0: float, width: float,
                         y: float, si: int) -> tuple[str, float]:
            heading = (sec.get("heading") or "").strip()
            items   = sec.get("items") or []
            # セクション見出しバーの色は最初のアイテムの色か、パレットから
            _hf, _hs = _STICKY_PALETTE[si % len(_STICKY_PALETTE)]
            per_row = max(1, int((width - SEC_PAD * 2 + GAP) / (STICKY_W + GAP)))
            n_rows  = (len(items) + per_row - 1) // max(per_row, 1) if items else 0
            items_h = n_rows * (STICKY_H + GAP) - (GAP if n_rows else 0)
            sec_h   = SEC_PAD + (SEC_HDR_H if heading else 0) + items_h + SEC_PAD
            bg  = (f'<rect x="{x0:.1f}" y="{y:.1f}" width="{width:.1f}"'
                   f' height="{sec_h}" rx="10" fill="#F8F8F8"'
                   f' stroke="#E0E0E0" stroke-width="1"/>\n')
            hdr = ""
            if heading:
                bg  += (f'<rect x="{x0:.1f}" y="{y:.1f}" width="{width:.1f}"'
                        f' height="{SEC_HDR_H}" rx="10" fill="{_hs}" opacity="0.18"/>\n')
                hdr  = (f'<text x="{x0 + 14}" y="{y + SEC_HDR_H * 0.68:.1f}"'
                        f' font-size="13" font-weight="bold" fill="{_hs}"'
                        f' font-family="system-ui,-apple-system,sans-serif">{esc(heading)}</text>\n')
            items_y = y + SEC_PAD + (SEC_HDR_H if heading else 0)
            stickies = ""
            for i, item in enumerate(items):
                row     = i // per_row
                col     = i % per_row
                cx      = x0 + SEC_PAD + col * (STICKY_W + GAP) + STICKY_W / 2
                cy      = items_y + row * (STICKY_H + GAP) + STICKY_H / 2
                f, s    = item_colors(item, si)   # 実物色を優先
                stickies += svg_sticky(cx, cy, item.get("text", ""),
                                       f, s, f"item-{si}-{i}",
                                       text_color=pen_color(item))
            return f'<g id="section-{si}">\n{bg}{hdr}{stickies}</g>\n', sec_h

        col2_secs  = [b for b in blocks if b.get("type") == "section" and b.get("column") == 2]
        has_two    = bool(col2_secs)
        cur_y      = PAGE_PAD + 56

        if has_two:
            col1_list = [b for b in blocks if b.get("type") == "section" and b.get("column", 1) != 2]
            col2_list  = col2_secs
            col_w      = (INNER_W - SEC_GAP) / 2
            left_y = right_y = cur_y
            for sec in col1_list:
                s, h = grid_section(sec, PAGE_PAD, col_w, left_y, sec_idx)
                body_svg += s; left_y += h + SEC_GAP; sec_idx += 1
            for sec in col2_list:
                s, h = grid_section(sec, PAGE_PAD + col_w + SEC_GAP, col_w, right_y, sec_idx)
                body_svg += s; right_y += h + SEC_GAP; sec_idx += 1
            cur_y = max(left_y, right_y)
        else:
            for block in blocks:
                if block.get("type") == "section":
                    s, h = grid_section(block, PAGE_PAD, INNER_W, cur_y, sec_idx)
                    body_svg += s; cur_y += h + SEC_GAP; sec_idx += 1
                elif block.get("type") == "table":
                    s, h = svg_table(block, PAGE_PAD, cur_y, INNER_W, f"table-{tbl_idx}")
                    body_svg += s; cur_y += h + SEC_GAP; tbl_idx += 1

        total_canvas_h = cur_y + PAGE_PAD

    # ── タイトル ──
    title_svg = ""
    if title:
        title_svg = (f'<text x="{PAGE_PAD}" y="{PAGE_PAD + 38}" font-size="22"'
                     f' font-weight="bold" fill="#1E1B4B"'
                     f' font-family="system-ui,-apple-system,sans-serif">{esc(title)}</text>\n')

    # ── 矢じりマーカー定義（各ペン色ぶん）──
    _ARROW_COLORS = list(PEN_COLOR_MAP.values()) + ["#1A1A1A"]
    _marker_defs  = "\n".join(
        f'  <marker id="ah-{hex_col.lstrip("#")}" markerWidth="8" markerHeight="6"'
        f' refX="7" refY="3" orient="auto" markerUnits="strokeWidth">'
        f'<path d="M0,0 L0,6 L8,3 z" fill="{hex_col}"/></marker>'
        for hex_col in dict.fromkeys(_ARROW_COLORS)  # deduplicate
    )

    svg_doc = f"""<svg width="{CANVAS_W}" height="{total_canvas_h}" xmlns="http://www.w3.org/2000/svg">
<defs>
  <filter id="dropshadow" x="-10%" y="-10%" width="120%" height="120%">
    <feDropShadow dx="2" dy="2" stdDeviation="3" flood-opacity="0.12"/>
  </filter>
{_marker_defs}
</defs>
<rect width="{CANVAS_W}" height="{total_canvas_h}" fill="#FFFFFF"/>
<g id="title">{title_svg}</g>
{body_svg}
</svg>"""
    return svg_doc.encode("utf-8")


def generate_html(data: dict) -> bytes:
    """iPhone Safari でそのまま開けるHTMLを sections ベースで生成"""
    title = data.get("title", "メモ")
    CMAP = {
        "black": "#1E1B4B", "white": "#6B7280", "red": "#DC2626",
        "blue": "#2563EB", "green": "#16A34A", "yellow": "#B45309",
        "orange": "#EA580C", "purple": "#7C3AED", "pink": "#DB2777",
        "gray": "#6B7280", "brown": "#92400E",
    }
    # 付箋の背景色マップ (fill, border)
    BG_HTML_MAP: dict[str, tuple[str, str]] = {
        "yellow": ("#FFF9C4", "#F9A825"),
        "pink":   ("#FCE4EC", "#E91E63"),
        "green":  ("#C8E6C9", "#2E7D32"),
        "blue":   ("#E3F2FD", "#1565C0"),
        "orange": ("#FFF3E0", "#E65100"),
        "purple": ("#F3E5F5", "#6A1B9A"),
        "white":  ("#FFFFFF", "#9CA3AF"),
    }

    def render_item(item):
        etype    = item.get("type", "text")
        content  = item.get("text", "")
        shape    = item.get("shape", "none")
        color    = CMAP.get(item.get("color", "black"), "#1E1B4B")
        bold     = item.get("bold", False) or etype == "heading"
        fw       = "700" if bold else "400"
        bg_color = item.get("bg_color", "none") or "none"
        is_bullet = etype == "bullet"

        # 付箋スタイル（bg_color あり）
        if bg_color != "none" and bg_color in BG_HTML_MAP:
            bg_fill, bd_color = BG_HTML_MAP[bg_color]
            inner = (f"<span class='sticky' style='background:{bg_fill};"
                     f" border-color:{bd_color}; font-weight:{fw};'>{content}</span>")
            return f"<li style='list-style:none; padding-left:0;'>{inner}</li>" if is_bullet else inner

        # 図形バッジ（rect / ellipse）
        if shape == "rect":
            inner = (f"<span class='badge rect' style='color:{color};"
                     f" border-color:{color}; font-weight:{fw};'>{content}</span>")
            return f"<li style='list-style:none; padding-left:0;'>{inner}</li>" if is_bullet else inner
        if shape == "ellipse":
            inner = (f"<span class='badge ellipse' style='color:{color};"
                     f" border-color:{color}; font-weight:{fw};'>{content}</span>")
            return f"<li style='list-style:none; padding-left:0;'>{inner}</li>" if is_bullet else inner

        if etype == "heading":
            return f"<h3 style='color:{color};'>{content}</h3>"
        if etype == "arrow":
            return f"<p class='arrow'>&rarr; {content}</p>"
        if is_bullet:
            return f"<li style='color:{color}; font-weight:{fw};'>{content}</li>"
        return f"<p style='color:{color}; font-weight:{fw};'>{content}</p>"

    def render_section(sec):
        heading = sec.get("heading", "").strip()
        items   = sec.get("items") or []
        out = ""
        if heading:
            out += f"<h2 class='sec-heading'>{heading}</h2>"
        # セクション見出しと同じテキストの先頭 heading アイテムは重複なのでスキップ
        start_idx = 0
        if (items and heading
                and items[0].get("type") == "heading"
                and items[0].get("text", "").strip() == heading):
            start_idx = 1
        # 箇条書きをまとめて <ul> に（インデックスを使わず状態フラグで管理）
        in_ul = False
        for item in items[start_idx:]:
            if item.get("type") == "bullet":
                if not in_ul:
                    out += "<ul style='padding-left:0; margin:0.4rem 0;'>"
                    in_ul = True
                out += render_item(item)
            else:
                if in_ul:
                    out += "</ul>"
                    in_ul = False
                out += render_item(item)
        if in_ul:
            out += "</ul>"
        return out

    def render_table_block(tbl: dict) -> str:
        headers = tbl.get("headers") or []
        trows   = tbl.get("rows") or []
        if not headers and not trows:
            return ""
        out = "<table>"
        if headers:
            out += "<tr>" + "".join(f"<th>{h}</th>" for h in headers) + "</tr>"
        for i, row in enumerate(trows):
            cls = "even" if i % 2 == 0 else ""
            out += f"<tr class='{cls}'>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>"
        out += "</table>"
        return out

    blocks = data.get("blocks") or []
    col2_exists = any(b.get("column", 1) == 2
                      for b in blocks if b.get("type") == "section")

    if col2_exists:
        # 2列モード：セクションを左右に分けて並列、表はその後に出現順
        col1_secs = [b for b in blocks if b.get("type") == "section" and b.get("column", 1) != 2]
        col2_secs = [b for b in blocks if b.get("type") == "section" and b.get("column", 1) == 2]
        body  = "<div class='two-col'>"
        body += "<div class='col'>" + "".join(render_section(s) for s in col1_secs) + "</div>"
        body += "<div class='col'>" + "".join(render_section(s) for s in col2_secs) + "</div>"
        body += "</div>"
        for b in blocks:
            if b.get("type") == "table":
                body += render_table_block(b)
    else:
        # 1列モード：blocks 出現順にそのまま描画
        body = ""
        for b in blocks:
            if b.get("type") == "section":
                body += render_section(b)
            elif b.get("type") == "table":
                body += render_table_block(b)

    ts = datetime.now(JST).strftime("%Y年%m月%d日 %H:%M")
    html = f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title} — パシャッと</title>
<style>
  * {{ box-sizing: border-box; }}
  body {{ font-family: -apple-system, 'Hiragino Sans', 'Yu Gothic', sans-serif;
          background: #F5F3FF; margin: 0; padding: 1rem; color: #1E1B4B;
          font-size: 16px; line-height: 1.7; }}
  .card {{ background: #fff; border-radius: 16px; padding: 1.5rem 1.6rem;
           box-shadow: 0 2px 16px rgba(109,40,217,0.12); max-width: 720px;
           margin: 0 auto; }}
  .header {{ background: linear-gradient(135deg,#6D28D9,#A855F7,#EC4899);
             border-radius: 12px; padding: 1rem 1.4rem; margin-bottom: 1.2rem; }}
  .header h1 {{ color:#fff; margin:0; font-size:1.3rem; }}
  .header p  {{ color:rgba(255,255,255,0.75); margin:0.25rem 0 0; font-size:0.8rem; }}
  hr {{ border:none; border-top:1px solid #EDE9FE; margin:1rem 0; }}
  h2.sec-heading {{ font-size:1.05rem; font-weight:800; color:#4C1D95;
                    border-left:4px solid #7C3AED; padding-left:0.6rem;
                    margin:1.2rem 0 0.5rem; }}
  h3 {{ font-size:1rem; margin:0.8rem 0 0.3rem; }}
  ul {{ padding-left:1.4rem; margin:0.3rem 0; }}
  li {{ margin:0.2rem 0; }}
  p  {{ margin:0.3rem 0; }}
  .arrow {{ color:#7C3AED; font-weight:600; }}
  .badge {{ display:inline-block; border:2px solid; border-radius:6px;
            padding:0.15rem 0.6rem; margin:0.25rem 0; font-weight:600; }}
  .badge.ellipse {{ border-radius:50px; padding:0.15rem 0.9rem; }}
  .sticky {{ display:inline-block; border:2px solid; border-radius:8px;
             padding:0.3rem 0.75rem; margin:0.25rem 0; font-size:0.95rem;
             box-shadow:1px 2px 5px rgba(0,0,0,0.08); line-height:1.5; }}
  .two-col {{ display:flex; gap:1.2rem; }}
  .col {{ flex:1; min-width:0; }}
  table {{ border-collapse:collapse; width:100%; margin:1rem 0; font-size:0.9rem; }}
  th {{ background:#7C3AED; color:#fff; padding:0.4rem 0.6rem;
        text-align:left; border:1px solid #C4B5FD; }}
  td {{ padding:0.35rem 0.6rem; border:1px solid #EDE9FE; }}
  tr.even td {{ background:#F5F3FF; }}
  footer {{ text-align:center; color:#A78BFA; font-size:0.75rem; margin-top:1.5rem; }}
  @media (max-width:520px) {{ .two-col {{ flex-direction:column; }} }}
</style>
</head>
<body>
<div class="card">
  <div class="header">
    <h1>{title}</h1>
    <p>{ts} 変換 | パシャッと</p>
  </div>
  <hr>
  {body}
</div>
<footer>パシャッと — AI文字認識による自動変換</footer>
</body>
</html>"""
    return html.encode("utf-8")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ─── 画面表示 ────────────────────────────────────────────
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

st.markdown("""
<div class='top-bar'>
  <div class='top-bar-logo'>パシャッ<em>と</em></div>
  <div class='top-bar-tagline'>メモをスライドに変換</div>
</div>
""", unsafe_allow_html=True)

# ── APIキー入力（未設定時のみサイドバーに表示） ──
if not get_api_key():
    with st.sidebar:
        st.markdown("### 管理者設定")
        st.markdown("AIを利用するためのキーを入力してください。")
        key_in = st.text_input(
            "認証キー",
            type="password",
            placeholder="sk-ant-...",
            help="管理者から受け取ったキーを入力してください。",
        )
        if key_in:
            st.session_state["api_key_input"] = key_in
            st.success("設定できました")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ステップ１　写真を選ぶ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.markdown(
    "<div class='step'><span class='snum'>１</span>ファイルを選ぶ</div>",
    unsafe_allow_html=True)
st.markdown(
    "<div class='hint'>メモ・ホワイトボード・資料の写真やPDFをAIが読み取ってスライドに変換します。</div>",
    unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    "画像ファイルを選ぶ",
    type=["jpg", "jpeg", "png", "webp", "pdf"],
    label_visibility="collapsed",
    help="JPG・PNG・WebP・PDF 形式に対応しています。",
)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 画像プレビュー ＆ 前処理
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# session_state から前回の処理済みデータを取得（ファイル未選択時のフォールバック）
img_data    = st.session_state.get("_img_data")
media_type  = st.session_state.get("_media_type", "image/jpeg")
is_portrait = st.session_state.get("_is_portrait", False)

if uploaded_file:
    st.markdown("---")

    is_pdf = uploaded_file.name.lower().endswith(".pdf")

    # ── PDF：ページ選択 ──
    pdf_page_num = 0
    if is_pdf:
        try:
            raw_bytes  = uploaded_file.getvalue()
            page_count = get_pdf_info(raw_bytes)
            if page_count > 1:
                st.markdown(
                    f"<div class='hint'>このPDFは <strong>{page_count} ページ</strong> あります。"
                    "変換するページを選んでください。</div>",
                    unsafe_allow_html=True)
                pdf_page_num = st.number_input(
                    "ページ番号",
                    min_value=1, max_value=page_count, value=1, step=1,
                    help=f"1〜{page_count} の範囲で選べます",
                ) - 1   # 0-indexed に変換
            else:
                st.caption("PDF（1ページ）")
        except Exception as pdf_info_err:
            st.error(f"PDFの読み込みに失敗しました。（{pdf_info_err}）")

    do_trim = st.toggle(
        "余白を自動カット（おすすめ）",
        value=True,
        help="不要な余白を取り除いて、読み取り精度をアップします。",
    )
    st.caption("明るさ・コントラストはAIが自動調整")

    try:
        file_bytes = uploaded_file.getvalue()

        # PDF の場合は先にページを画像に変換してから処理
        if is_pdf:
            pil_from_pdf = pdf_page_to_pil(file_bytes, page_num=pdf_page_num)
            buf_for_cache = io.BytesIO()
            pil_from_pdf.save(buf_for_cache, format="JPEG", quality=95)
            proc_bytes = buf_for_cache.getvalue()
        else:
            proc_bytes = file_bytes

        # キャッシュキー：ファイル名＋サイズ（＋PDFならページ番号も含める）
        file_id = f"{uploaded_file.name}_{len(file_bytes)}_p{pdf_page_num}"

        # キャッシュ済み処理（同一ファイル＋設定ならリランしても即返却）
        img_data, media_type, is_portrait, preview_bytes = cached_process_image(
            proc_bytes, do_trim
        )

        # ファイルが変わったら前の変換結果・編集内容をリセット
        if st.session_state.get("_file_id") != file_id:
            for k in ("extracted", "pptx_bytes", "html_bytes"):
                st.session_state.pop(k, None)
            for k in [k for k in st.session_state if k.startswith("edit_")]:
                del st.session_state[k]
            st.session_state["_file_id"] = file_id

        # session_state に保持（リラン後もボタン押下で参照できるよう）
        st.session_state["_img_data"]       = img_data
        st.session_state["_media_type"]     = media_type
        st.session_state["_is_portrait"]    = is_portrait
        st.session_state["_preview_bytes"]  = preview_bytes

        orient_label = "縦" if is_portrait else "横"
        cap_label    = f"ページ {pdf_page_num + 1}" if is_pdf else "変換元"
        st.caption(f"スライド向き：{orient_label}で出力")
        st.image(Image.open(io.BytesIO(preview_bytes)),
                 caption=cap_label, use_container_width=True)

    except Exception as img_err:
        st.error(f"ファイルの読み込みに失敗しました。別のファイルをお試しください。\n（{img_err}）")
        img_data = None

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ステップ２　変換する
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
st.markdown("---")
st.markdown(
    "<div class='step'><span class='snum'>２</span>変換する</div>",
    unsafe_allow_html=True)

if not uploaded_file:
    st.markdown(
        "<div class='hint' style='border-color:#C4B5FD; background:#F5F3FF;'>"
        "まずファイルを選んでください。"
        "</div>",
        unsafe_allow_html=True)

# ── モデル選択（コスト vs 精度） ──
MODEL_OPTIONS = {
    "高精度（おすすめ）　約4円/回": "claude-sonnet-4-6",
    "省エネ・高速　約1円/回": "claude-haiku-4-5-20251001",
}
selected_model_label = st.radio(
    "読み取りモード",
    list(MODEL_OPTIONS.keys()),
    index=0,
    horizontal=True,
    help="「省エネ」は速くて安い分、複雑な手書きや細かい文字の読み取り精度が下がる場合があります。",
)
selected_model = MODEL_OPTIONS[selected_model_label]

convert_btn = st.button(
    "スライドに変換する",
    type="primary",
    use_container_width=True,
    disabled=not uploaded_file,
)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 変換処理
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
if convert_btn and not img_data:
    st.error("写真を正しく読み込めませんでした。別の写真を選んでもう一度お試しください。")

if convert_btn and img_data:
    try:
        api_key = get_api_key()
        if not api_key:
            st.warning("認証キーが設定されていません。左上のメニューから設定してください。")
            st.stop()

        # ── ① 準備（瞬時に完了）──
        prog = st.progress(0, text="写真を準備しています…")
        prog.progress(30, text="写真を準備しています…")

        # ── ② AI読み取り（最も時間がかかる → スピナーで安心感を演出）──
        with st.spinner("AIが読み取り中です…"):
            extracted = analyze_with_claude(img_data, media_type, api_key,
                                            model=selected_model)

        # ── ③ スライド生成（数秒）──
        prog.progress(80, text="スライドを作成しています…")

        is_portrait = st.session_state.get("_is_portrait", False)
        st.session_state["extracted"]   = extracted
        st.session_state["pptx_bytes"]  = generate_pptx(extracted, is_portrait=is_portrait)
        st.session_state["html_bytes"]  = generate_html(extracted)
        # 新規変換時は編集キーをクリア（前回の編集内容を引き継がせない）
        for k in [k for k in st.session_state if k.startswith("edit_")]:
            del st.session_state[k]

        prog.progress(100, text="完了しました！")
        prog.empty()
        st.success("変換できました！")

    except Exception as e:
        err_msg = str(e)
        if "api_key" in err_msg.lower() or "authentication" in err_msg.lower() or "401" in err_msg:
            st.error("認証キーが正しくありません。左上のメニューから確認してください。")
        elif "json" in err_msg.lower() or "JSON" in err_msg or "応答" in err_msg:
            st.error(f"文字の読み取り結果を処理できませんでした。もう一度お試しください。\n（詳細：{err_msg[:200]}）")
        elif "timeout" in err_msg.lower() or "connection" in err_msg.lower():
            st.error("通信エラーが発生しました。インターネット接続をご確認のうえ、もう一度お試しください。")
        elif "overloaded" in err_msg.lower() or "529" in err_msg:
            st.error("AIサーバーが混み合っています。しばらく待ってから、もう一度お試しください。")
        else:
            st.error(f"エラーが発生しました。もう一度お試しください。\n（詳細：{err_msg}）")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ステップ３　保存する
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
if "extracted" in st.session_state:
    extracted   = st.session_state["extracted"]
    pptx_bytes  = st.session_state.get("pptx_bytes")
    html_bytes  = st.session_state.get("html_bytes")

    st.markdown("---")
    st.markdown(
        "<div class='step'><span class='snum'>３</span>保存する</div>",
        unsafe_allow_html=True)

    elements = extracted.get("elements", [])
    el_count = len(elements)
    st.markdown(
        f"<div class='done-box'>完了 — 読み取り：<strong>{el_count} 件</strong></div>",
        unsafe_allow_html=True)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # テキスト編集エリア（元画像を見ながら修正）
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    st.markdown("---")
    st.markdown(
        "<div class='step'><span class='snum' style='font-size:0.95rem;'>✏</span>"
        "内容を修正する（任意）</div>",
        unsafe_allow_html=True)
    st.markdown(
        "<div class='hint'>読み取り結果を修正できます。<strong>[?]</strong> は読み取りが不確かな箇所です。</div>",
        unsafe_allow_html=True)

    # ── 元画像プレビュー ──
    preview_bytes_sv = st.session_state.get("_preview_bytes")
    if preview_bytes_sv:
        st.image(Image.open(io.BytesIO(preview_bytes_sv)),
                 caption="変換元", use_container_width=True)

    # 編集フォーム（全幅 — スマホ・PCどちらでも使いやすい）
    TYPE_LABEL  = {"heading": "見出し", "bullet": "箇条書き",
                   "text": "本文", "arrow": "矢印"}
    SHAPE_LABEL = {"rect": "【四角】", "ellipse": "（丸）"}

    # タイトル
    title_key = "edit_title"
    if title_key not in st.session_state:
        st.session_state[title_key] = extracted.get("title", "")
    st.markdown(
        "<p style='font-size:0.88rem; font-weight:700; color:#1E1B4B;"
        " margin:0.5rem 0 0.1rem;'>タイトル</p>",
        unsafe_allow_html=True)
    st.text_input("タイトル", key=title_key, label_visibility="collapsed")

    # ── 精度チェックボタン ──
    _VF_BADGE = {
        "high":   ("🟢", "#16A34A", "正確"),
        "medium": ("🟡", "#CA8A04", "やや不確か"),
        "low":    ("🔴", "#DC2626", "誤読の可能性あり"),
    }
    if st.button("🔍 識字精度を確認する（約1円）",
                 type="secondary", use_container_width=True):
        api_key_v = get_api_key()
        img_data_v = st.session_state.get("_img_data")
        media_v    = st.session_state.get("_media_type", "image/jpeg")
        if not api_key_v or not img_data_v:
            st.warning("APIキーまたは画像データがありません。")
        else:
            with st.spinner("AIが識字精度を検証中です…"):
                try:
                    vr = verify_extraction(img_data_v, media_v, api_key_v, extracted)
                    st.session_state["verification"] = vr
                    low_cnt    = sum(1 for v in vr.values() if v["confidence"] == "low")
                    medium_cnt = sum(1 for v in vr.values() if v["confidence"] == "medium")
                    if low_cnt == 0 and medium_cnt == 0:
                        st.success("精度チェック完了：すべて高信頼度です。")
                    else:
                        st.warning(
                            f"精度チェック完了：🔴 要確認 {low_cnt}件　"
                            f"🟡 やや不確か {medium_cnt}件　"
                            f"→ 下の「再読み取り」ボタンで修正できます")
                    # 修正案がある項目を session_state に反映（上書き確認は UI 側で）
                    for (bi, ii), v in vr.items():
                        if v.get("correction"):
                            ekey = f"edit_b{bi}_i{ii}"
                            # 未編集の場合のみ自動反映
                            orig = ""
                            for bk in (extracted.get("blocks") or []):
                                if bk.get("type") == "section":
                                    pass
                            blocks_tmp = extracted.get("blocks") or []
                            if bi < len(blocks_tmp):
                                items_tmp = blocks_tmp[bi].get("items") or []
                                if ii < len(items_tmp):
                                    orig = items_tmp[ii].get("text", "")
                            if st.session_state.get(ekey, orig) == orig:
                                st.session_state[ekey] = v["correction"]
                except Exception as ex:
                    st.error(f"検証エラー: {ex}")

    verification = st.session_state.get("verification", {})

    # ── 該当箇所だけ高精度モデルで再読み取り ──
    uncertain_cnt = sum(1 for v in verification.values()
                        if v.get("confidence") in ("medium", "low"))
    if verification and uncertain_cnt > 0:
        rr_model_map = {
            f"Sonnet（高精度）　約1円": "claude-sonnet-4-6",
        }
        rr_model_lbl = st.selectbox(
            "再読み取りモデル", list(rr_model_map.keys()),
            label_visibility="collapsed")
        rr_model = rr_model_map[rr_model_lbl]
        if st.button(
                f"✨ 不確かな {uncertain_cnt} 件だけ高精度で再読み取り",
                type="primary", use_container_width=True):
            api_key_r  = get_api_key()
            img_data_r = st.session_state.get("_img_data")
            media_r    = st.session_state.get("_media_type", "image/jpeg")
            if not api_key_r or not img_data_r:
                st.warning("APIキーまたは画像データがありません。")
            else:
                with st.spinner(f"高精度モデルで {uncertain_cnt} 件を再読み取り中…"):
                    try:
                        corrections = reread_uncertain(
                            img_data_r, media_r, api_key_r,
                            extracted, verification, model=rr_model)
                        fixed = 0
                        for (bi, ii), new_text in corrections.items():
                            ekey = f"edit_b{bi}_i{ii}"
                            if new_text and st.session_state.get(ekey) != new_text:
                                st.session_state[ekey] = new_text
                                # 検証結果も high に更新
                                if (bi, ii) in st.session_state["verification"]:
                                    st.session_state["verification"][(bi, ii)][
                                        "confidence"] = "high"
                                fixed += 1
                        st.success(f"{fixed} 件を修正しました。"
                                   "「作り直す」で反映してください。")
                    except Exception as ex:
                        st.error(f"再読み取りエラー: {ex}")
        st.markdown("---")

    # ── blocks を出現順に編集フォームとして表示 ──
    edit_blocks = extracted.get("blocks") or []
    global_idx  = 0   # アイテム全体を通した連番（色ラベル用）
    tbl_count   = 0   # 表ブロックの通し番号

    for bi, block in enumerate(edit_blocks):
        btype = block.get("type")

        if btype == "section":
            sec_heading = (block.get("heading") or "").strip()
            if sec_heading:
                col_lbl = 1 if block.get("column", 1) != 2 else 2
                st.markdown(
                    f"<p style='margin:0.8rem 0 0.2rem; font-size:0.85rem;"
                    f" color:#7C3AED; font-weight:700;'>── {sec_heading}"
                    f"（列{col_lbl}）</p>",
                    unsafe_allow_html=True)
            for ii, item in enumerate(block.get("items") or []):
                etype = item.get("type", "text")
                shape = item.get("shape", "none")
                lbl   = SHAPE_LABEL.get(shape, "") + TYPE_LABEL.get(etype, "本文")
                chex  = item_color_hex(global_idx)
                ekey  = f"edit_b{bi}_i{ii}"
                if ekey not in st.session_state:
                    st.session_state[ekey] = item.get("text", "")

                # 精度バッジ（検証済みの場合）
                vinfo = verification.get((bi, ii))
                badge_html = ""
                if vinfo:
                    icon, col, tip = _VF_BADGE.get(
                        vinfo["confidence"], ("🟡", "#CA8A04", ""))
                    badge_html = (
                        f"<span style='font-size:0.78rem; color:{col};"
                        f" font-weight:700; margin-left:6px;'>{icon} {tip}</span>")

                st.markdown(
                    f"<p style='font-size:0.88rem; font-weight:700; color:#1E1B4B;"
                    f" margin:0.6rem 0 0.1rem;'>"
                    f"<span style='display:inline-block; width:13px; height:13px;"
                    f" background:{chex}; border-radius:3px; margin-right:5px;"
                    f" vertical-align:middle;'></span>{lbl}{badge_html}</p>",
                    unsafe_allow_html=True)
                st.text_input(lbl, key=ekey, label_visibility="collapsed")
                global_idx += 1

        elif btype == "table":
            tbl_count += 1
            headers = list(block.get("headers") or [])
            rows    = list(block.get("rows") or [])
            n_cols  = max(len(headers), max((len(r) for r in rows), default=0))
            if n_cols == 0:
                continue
            st.markdown(
                f"<p style='font-size:0.92rem; font-weight:800; color:#1E1B4B;"
                f" margin:1.1rem 0 0.2rem;'>📊 表{tbl_count}</p>",
                unsafe_allow_html=True)
            if headers:
                st.markdown(
                    "<p style='font-size:0.82rem; color:#6B7280;"
                    " margin:0.3rem 0 0.1rem;'>ヘッダー行</p>",
                    unsafe_allow_html=True)
                h_cols = st.columns(n_cols)
                for ci in range(n_cols):
                    hkey = f"edit_b{bi}_h{ci}"
                    if hkey not in st.session_state:
                        st.session_state[hkey] = headers[ci] if ci < len(headers) else ""
                    with h_cols[ci]:
                        st.text_input(f"H{ci+1}", key=hkey, label_visibility="collapsed")
            for ri, row in enumerate(rows):
                if ri == 0:
                    st.markdown(
                        "<p style='font-size:0.82rem; color:#6B7280;"
                        " margin:0.3rem 0 0.1rem;'>データ行</p>",
                        unsafe_allow_html=True)
                r_cols = st.columns(n_cols)
                for ci in range(n_cols):
                    rkey = f"edit_b{bi}_r{ri}_{ci}"
                    if rkey not in st.session_state:
                        st.session_state[rkey] = row[ci] if ci < len(row) else ""
                    with r_cols[ci]:
                        st.text_input(f"R{ri+1}C{ci+1}", key=rkey,
                                      label_visibility="collapsed")

    # ── 作り直しボタン ──
    if st.button("この内容でスライドを作り直す",
                 type="secondary", use_container_width=True):
        edited = copy.deepcopy(extracted)
        edited["title"] = st.session_state.get("edit_title", edited.get("title", ""))

        # blocks を session_state の値で上書き
        for bi, block in enumerate(edited.get("blocks") or []):
            if block.get("type") == "section":
                for ii, item in enumerate(block.get("items") or []):
                    ekey = f"edit_b{bi}_i{ii}"
                    if ekey in st.session_state:
                        item["text"] = st.session_state[ekey]
            elif block.get("type") == "table":
                headers = list(block.get("headers") or [])
                for ci in range(len(headers)):
                    hkey = f"edit_b{bi}_h{ci}"
                    if hkey in st.session_state:
                        block["headers"][ci] = st.session_state[hkey]
                for ri, row in enumerate(block.get("rows") or []):
                    for ci in range(len(row)):
                        rkey = f"edit_b{bi}_r{ri}_{ci}"
                        if rkey in st.session_state:
                            block["rows"][ri][ci] = st.session_state[rkey]

        ip = st.session_state.get("_is_portrait", False)
        st.session_state["extracted"]  = edited
        st.session_state["pptx_bytes"] = generate_pptx(edited, is_portrait=ip)
        st.session_state["html_bytes"] = generate_html(edited)
        st.success("スライドを更新しました。下のボタンから保存してください。")

    st.markdown("---")
    ts = datetime.now(JST).strftime("%Y%m%d_%H%M%S")

    # ── iPhone 向け HTML ダウンロード（メイン） ──
    if html_bytes:
        st.markdown(
            "<div style='font-size:0.9rem; color:#7C3AED; margin-bottom:0.4rem;"
            " font-weight:600;'>iPhone でそのまま開けます</div>",
            unsafe_allow_html=True)
        st.download_button(
            label="HTMLで保存する（iPhone対応）",
            data=html_bytes,
            file_name=f"パシャッと_{ts}.html",
            mime="text/html",
            use_container_width=True,
        )

    # ── PowerPoint ダウンロード（サブ） ──
    if pptx_bytes:
        st.markdown(
            "<div style='font-size:0.9rem; color:#6B7280; margin:0.8rem 0 0.4rem;"
            " font-weight:600;'>PC・PowerPoint 用</div>",
            unsafe_allow_html=True)
        st.download_button(
            label="PowerPoint で保存する",
            data=pptx_bytes,
            file_name=f"パシャッと_{ts}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

    # ── Figma 用 SVG ダウンロード ──
    if extracted:
        st.markdown(
            "<div style='font-size:0.9rem; color:#0EA5E9; margin:0.8rem 0 0.4rem;"
            " font-weight:600;'>Figma・イラストレーター用</div>",
            unsafe_allow_html=True)
        preview_bytes = st.session_state.get("_preview_bytes")
        svg_bytes = generate_svg(extracted, preview_bytes=preview_bytes)
        st.download_button(
            label="📌 ポストイット配置をSVGで保存する（Figma対応）",
            data=svg_bytes,
            file_name=f"パシャッと_{ts}.svg",
            mime="image/svg+xml",
            use_container_width=True,
        )


# ── フッター ──
st.markdown("<br>", unsafe_allow_html=True)
st.markdown(
    "<div style='text-align:center; color:#A78BFA; font-size:0.85rem; line-height:2;'>"
    "AI文字認識による自動変換サービス<br>&copy; 2026 パシャッと"
    "</div>",
    unsafe_allow_html=True)
